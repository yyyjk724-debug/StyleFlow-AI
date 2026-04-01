import streamlit as st
import pandas as pd
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

# 1. 페이지 설정
st.set_page_config(page_title="StyleFlow AI", page_icon="🌿", layout="wide")

# 2. 다크모드 지원 및 고급 디자인 CSS
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Pretendard:wght@400;600;900&display=swap');
    
    /* 기본 테마 설정 (라이트/다크 대응) */
    .main { font-family: 'Pretendard', sans-serif; transition: all 0.3s; }
    
    /* 제목 스타일: 아주 크게, 굵게 900 */
    .main-title {
        font-size: 4.2rem;
        font-weight: 900;
        letter-spacing: -3px;
        margin-bottom: 0.5rem;
        text-align: center;
        color: var(--title-color);
    }
    
    /* 다크모드/라이트모드 변수 설정 */
    [data-theme="light"] { --title-color: #111827; --bg-color: #f9fafb; --text-color: #374151; }
    [data-theme="dark"] { --title-color: #ffffff; --bg-color: #111827; --text-color: #e5e7eb; }

    /* 입력/결과 박스 애니메이션 */
    .stTextArea textarea {
        border-radius: 16px;
        border: 1px solid #d1d5db;
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
        font-size: 16px !important;
        padding: 20px;
    }
    .stTextArea textarea:focus {
        border-color: #10b981 !important;
        transform: translateY(-4px);
        box-shadow: 0 10px 20px rgba(16, 185, 129, 0.1) !important;
    }

    /* 표 섹션 디자인 */
    .table-section {
        background-color: rgba(16, 185, 129, 0.05);
        border-radius: 16px;
        padding: 25px;
        border: 1px solid #10b981;
        margin-top: 25px;
        animation: slideUp 0.5s ease-out;
    }
    
    @keyframes slideUp { from { opacity: 0; transform: translateY(15px); } to { opacity: 1; transform: translateY(0); } }
    </style>
    """, unsafe_allow_html=True)

# 3. 사이드바 - 흑백(다크) 모드 스위치
with st.sidebar:
    st.header("🌓 테마 설정")
    dark_mode = st.toggle("흑백(다크) 모드 활성화")
    st.divider()
    st.caption("밤에 작업할 때 눈을 보호하세요.")

# 테마에 따른 배경색 강제 적용
if dark_mode:
    st.markdown('<style>.main { background-color: #111827 !important; color: white !important; }</style>', unsafe_allow_html=True)
else:
    st.markdown('<style>.main { background-color: #f9fafb !important; color: #374151 !important; }</style>', unsafe_allow_html=True)

# 4. 헤더
st.markdown('<h1 class="main-title">StyleFlow AI</h1>', unsafe_allow_html=True)
st.markdown('<p style="text-align: center; font-size: 1.2rem; opacity: 0.8;">서식 세탁부터 PPT 표 생성까지, 단 한 번의 복붙으로.</p>', unsafe_allow_html=True)

# 5. 좌우 메인 레이아웃 (실시간 변환)
col_in, col_out = st.columns(2, gap="large")

with col_in:
    st.markdown("### 📥 AI 답변 붙여넣기")
    user_content = st.text_area("Input", height=500, placeholder="ChatGPT/Gemini 답변을 넣으세요...", label_visibility="collapsed")

with col_out:
    st.markdown("### ✅ 변환 결과")
    if user_content:
        # 서식 정제 로직
        clean_text = re.sub(r'(^|\n)[*#>-]\s?', r'\1', user_content)
        clean_text = re.sub(r'[*_~`]', '', clean_text)
        
        # 버튼 없이 실시간으로 결과창에 노출
        st.text_area("Output", value=clean_text, height=350, label_visibility="collapsed")
        st.markdown("<p style='color: #10b981; font-weight: 700;'>✔ Ctrl+A -> Ctrl+C 하여 PPT에 바로 붙여넣으세요!</p>", unsafe_allow_html=True)

        # 표 데이터 감지 시 전용 섹션
        if '|' in user_content:
            st.markdown('<div class="table-section">', unsafe_allow_html=True)
            st.markdown("<h4 style='color: #10b981;'>📊 표 데이터 자동 변환</h4>", unsafe_allow_html=True)
            try:
                lines = [l.strip() for l in user_content.split('\n') if '|' in l]
                valid_lines = [l for l in lines if '---' not in l]
                if len(valid_lines) > 1:
                    headers = [h.strip() for h in valid_lines[0].split('|') if h.strip()]
                    data = [[cell.strip() for cell in l.split('|') if cell.strip()] for l in valid_lines[1:]]
                    df = pd.DataFrame(data, columns=headers)
                    
                    # [추가 기능] PPT 표 자동 생성 엔진
                    prs = Presentation()
                    slide = prs.slides.add_slide(prs.slide_layouts[5]) # 빈 슬라이드
                    rows, cols = len(df) + 1, len(df.columns)
                    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(rows*0.5)).table
                    
                    for i, col_name in enumerate(df.columns): table.cell(0, i).text = col_name
                    for r in range(len(df)):
                        for c in range(len(df.columns)): table.cell(r+1, c).text = str(df.iloc[r, c])
                    
                    ppt_out = BytesIO()
                    prs.save(ppt_out)
                    
                    st.success("표가 포함된 PPT 파일을 준비했습니다!")
                    st.download_button("📂 표 전용 PPT 파일 다운로드", data=ppt_out.getvalue(), file_name="styleflow_table.pptx")
                    st.dataframe(df, use_container_width=True)
            except: pass
            st.markdown('</div>', unsafe_allow_html=True)
    else:
        st.info("왼쪽에 내용을 입력하면 실시간으로 서식이 제거됩니다.")

st.markdown("<br><hr><p style='text-align: center; opacity: 0.5; font-size: 0.8rem;'>© 2026 StyleFlow AI. Built for Seo-yeon's Academic Success.</p>", unsafe_allow_html=True)
